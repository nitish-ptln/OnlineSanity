"""
Telephony Manager - Demo Presentation Generator
Generates a professional PPT with dark theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Colors ───────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0F, 0x11, 0x1A)
BG_CARD     = RGBColor(0x1A, 0x1D, 0x2E)
BG_CARD2    = RGBColor(0x14, 0x16, 0x22)
ACCENT      = RGBColor(0x6C, 0x5C, 0xE7)  # Purple accent
ACCENT2     = RGBColor(0x00, 0xD2, 0xFF)  # Cyan accent
GREEN       = RGBColor(0x00, 0xE6, 0x76)
RED         = RGBColor(0xFF, 0x44, 0x44)
ORANGE      = RGBColor(0xFF, 0xA5, 0x02)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xBB, 0xBB, 0xCC)
MUTED       = RGBColor(0x88, 0x88, 0xAA)
YELLOW      = RGBColor(0xFF, 0xD9, 0x3D)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=LIGHT_GRAY):
    """Add textbox with multiple styled lines. Each line is (text, size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Segoe UI'
        p.space_after = Pt(6)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, size=15, color=LIGHT_GRAY, bullet_color=ACCENT2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Add bullet
        run1 = p.add_run()
        run1.text = "●  "
        run1.font.size = Pt(size - 2)
        run1.font.color.rgb = bullet_color
        run1.font.name = 'Segoe UI'
        
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(size)
        run2.font.color.rgb = color
        run2.font.name = 'Segoe UI'
        p.space_after = Pt(8)
    return txBox


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
set_bg(slide)

# Accent line at top
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

# Title area
add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
         "TELEPHONY MANAGER", 52, WHITE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.8),
         "Online Sanity Testing Platform", 30, ACCENT2, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(1.2),
         "A modern web-based tool for remote telephony sanity testing — replacing heavyweight\n"
         "Remote PC sessions with a lightweight, multi-user, real-time browser experience.",
         16, MUTED, False, PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, Inches(0), Inches(6.5), Inches(13.33), Inches(1), BG_CARD2)
add_text(slide, Inches(1), Inches(6.65), Inches(5), Inches(0.6),
         "LG Electronics  |  Connected Service Unit", 13, MUTED, False)
add_text(slide, Inches(7), Inches(6.65), Inches(5.5), Inches(0.6),
         "February 2026  |  v1.0", 13, MUTED, False, PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: Problem Statement — Why Not Remote PC?
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "THE PROBLEM  —  Remote PC Limitations", 30, WHITE, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

# Left Card - Remote PC Limitations
card = add_shape(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5), BG_CARD, RGBColor(0x33, 0x33, 0x44))
add_text(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
         "❌  Remote PC  (Current Method)", 20, RED, True)

limitations = [
    "Only ONE user can connect at a time (exclusive session)",
    "Requires full desktop streaming — high bandwidth needed",
    "Slow & laggy UI — screen refresh delays up to 2-3 seconds",
    "If connection drops, session is lost entirely",
    "No parallel testing — team waits for one person to finish",
    "Device physical access required to troubleshoot",
    "No automated regression — everything is manual",
    "DLT logs require separate setup on each PC",
    "No real-time device status dashboard",
    "No test history / export capability",
    "IP conflicts when multiple users try to access"
]
add_bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.2), Inches(4.5),
                limitations, 13, LIGHT_GRAY, RED)

# Right Card - Our Solution
card = add_shape(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.5), BG_CARD, RGBColor(0x22, 0x55, 0x44))
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.5),
         "✅  Telephony Manager  (Our Solution)", 20, GREEN, True)

solutions = [
    "Multi-user access — unlimited tabs, zero conflicts",
    "Lightweight browser UI — no screen streaming needed",
    "Instant command execution — sub-second response",
    "Session persists even if browser tab is refreshed",
    "Parallel testing — team works simultaneously",
    "Remote device access via ADB over network",
    "Built-in automated regression with CSV export",
    "DLT bridge auto-configures per user session",
    "Live device dashboard (IMEI, SIM, Radio, Region)",
    "Full test history with PASS/FAIL export",
    "Per-user isolation — no IP or port conflicts"
]
add_bullet_list(slide, Inches(7.1), Inches(2.15), Inches(5.5), Inches(4.5),
                solutions, 13, LIGHT_GRAY, GREEN)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: Key Features Overview
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "KEY FEATURES", 30, WHITE, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

features = [
    ("🖥️", "Live Device Dashboard", "Real-time status bar showing IMEI, SIM State,\nService State, Radio State, Region, SW Version", ACCENT2),
    ("⚡", "One-Click Commands", "Execute 50+ telephony commands with pass/fail\nvalidation using pattern matching", YELLOW),
    ("🔄", "Automated Regression", "Configure test sequences, set iterations, pause/\nresume, and export CSV reports", GREEN),
    ("📡", "DLT Bridge", "Auto-configures TCP proxy per user for DLT log\nforwarding — no manual setup needed", ACCENT),
    ("👥", "Multi-User Support", "Each user gets unique client ID, isolated config,\nseparate DLT ports — zero conflicts", ORANGE),
    ("🔔", "Global Notifications", "Real-time popups for ADB changes, IMEI writes,\nand system events across all users", RED),
    ("📊", "Multi-Model Support", "Toyota, JLR, BMW WAVE LOW/HIGH with auto-\ndetection based on device software version", ACCENT2),
    ("🛡️", "Admin Panel", "Add/remove models, commands, categories\ndynamically without code changes", YELLOW),
]

for i, (icon, title, desc, color) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.3 + row * 2.9)
    
    card = add_shape(slide, x, y, Inches(2.95), Inches(2.6), BG_CARD, RGBColor(0x33, 0x33, 0x44))
    add_text(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.6), Inches(0.5),
             icon, 32, color, False, PP_ALIGN.LEFT)
    add_text(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.6), Inches(0.45),
             title, 16, WHITE, True)
    add_text(slide, x + Inches(0.15), y + Inches(1.2), Inches(2.6), Inches(1.2),
             desc, 12, MUTED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: Technology Stack
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "TECHNOLOGY STACK", 30, WHITE, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

tech_stack = [
    ("FRONTEND", [
        ("HTML5 + CSS3 + JavaScript", "Modern responsive UI with dark theme"),
        ("Vanilla JS (No frameworks)", "Zero dependencies, fast load, easy maintenance"),
        ("CSS Variables & Glassmorphism", "Premium dark-mode design system"),
        ("Fetch API + Async/Await", "Non-blocking server communication"),
    ], ACCENT2),
    ("BACKEND", [
        ("Node.js + Express.js", "Lightweight, event-driven server"),
        ("ADB (Android Debug Bridge)", "Device communication via USB/TCP"),
        ("net (TCP Module)", "DLT proxy bridge for log forwarding"),
        ("child_process (execAsync)", "Shell command execution with timeout"),
    ], GREEN),
    ("TOOLING & DEPLOYMENT", [
        ("pkg (by Vercel)", "Compiles to standalone .exe — no Node.js needed"),
        ("Single Binary Distribution", "TelephonyManager.exe — just double-click to run"),
        ("sessionStorage + localStorage", "Client-side state persistence per tab"),
        ("JSON-based Config", "Dynamic model/command management via JSON files"),
    ], ORANGE),
]

for i, (category, items, color) in enumerate(tech_stack):
    x = Inches(0.4 + i * 4.2)
    card = add_shape(slide, x, Inches(1.4), Inches(4.0), Inches(5.6), BG_CARD, RGBColor(0x33, 0x33, 0x44))
    
    add_text(slide, x + Inches(0.2), Inches(1.5), Inches(3.5), Inches(0.5),
             category, 18, color, True)
    add_shape(slide, x + Inches(0.2), Inches(2.0), Inches(3.5), Inches(0.03), color)
    
    y_offset = 2.2
    for tech, desc in items:
        add_multi_text(slide, x + Inches(0.2), Inches(y_offset), Inches(3.5), Inches(0.9),
                      [(tech, 14, WHITE, True), (desc, 11, MUTED)])
        y_offset += 0.85


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: Architecture — How Server Connects
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "ARCHITECTURE  —  How the Server Connects", 30, WHITE, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

# Connection Flow Diagram
# User Browsers (left) → Server (center) → Devices (right)

# --- Users Column ---
add_text(slide, Inches(0.3), Inches(1.3), Inches(3), Inches(0.4),
         "BROWSER CLIENTS", 16, ACCENT2, True, PP_ALIGN.CENTER)

users = ["User A  (Tab 1)", "User B  (Tab 2)", "User C  (Tab 3)"]
for i, user in enumerate(users):
    y = Inches(1.9 + i * 1.2)
    card = add_shape(slide, Inches(0.3), y, Inches(3.0), Inches(0.9), BG_CARD, ACCENT2)
    add_text(slide, Inches(0.5), y + Inches(0.05), Inches(2.5), Inches(0.4),
             f"🌐  {user}", 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), y + Inches(0.45), Inches(2.5), Inches(0.35),
             f"clientId: client_{'abc'[i]}_{i+1}xxx", 10, MUTED, False, PP_ALIGN.CENTER)

# --- Arrows Left → Center ---
for i in range(3):
    y = Inches(2.2 + i * 1.2)
    add_text(slide, Inches(3.3), y, Inches(0.8), Inches(0.4),
             "──►", 18, ACCENT2, True, PP_ALIGN.CENTER)

# --- Server Column ---
add_text(slide, Inches(4.1), Inches(1.3), Inches(5), Inches(0.4),
         "EXPRESS.JS SERVER  (Port 3000)", 16, GREEN, True, PP_ALIGN.CENTER)

server_card = add_shape(slide, Inches(4.1), Inches(1.9), Inches(5.0), Inches(4.8), BG_CARD, GREEN)

server_items = [
    ("⚙️  Request Router", "Maps X-Client-ID header to per-user config\n(userConfigs Map: clientId → {serial, adb, model})", WHITE),
    ("📋  Command Engine", "Executes ADB commands via child_process\nwith timeout, caching & pattern validation", WHITE),
    ("💾  Device Cache", "10s TTL cache with Promise.allSettled\nStale-cache fallback prevents UI flickering", WHITE),
    ("🔔  Notification Bus", "Global notification array polled by all clients\nADB changes, IMEI writes broadcast to everyone", WHITE),
]

for i, (title, desc, color) in enumerate(server_items):
    y = Inches(2.0 + i * 1.1)
    add_text(slide, Inches(4.3), y, Inches(4.5), Inches(0.35),
             title, 13, color, True)
    add_text(slide, Inches(4.3), y + Inches(0.3), Inches(4.5), Inches(0.6),
             desc, 10, MUTED)

# --- Arrows Center → Right ---
for i in range(2):
    y = Inches(2.2 + i * 2.0)
    add_text(slide, Inches(9.1), y, Inches(0.8), Inches(0.4),
             "──►", 18, GREEN, True, PP_ALIGN.CENTER)

# --- Devices Column ---
add_text(slide, Inches(9.9), Inches(1.3), Inches(3), Inches(0.4),
         "TARGET DEVICES", 16, ORANGE, True, PP_ALIGN.CENTER)

devices = [
    ("📱  Device A (Toyota)", "adb1 -s SERIAL_A"),
    ("📱  Device B (BMW)", "adb1 -s SERIAL_B"),
]
for i, (dev, cmd) in enumerate(devices):
    y = Inches(1.9 + i * 2.0)
    card = add_shape(slide, Inches(9.9), y, Inches(3.0), Inches(1.4), BG_CARD, ORANGE)
    add_text(slide, Inches(10.1), y + Inches(0.1), Inches(2.6), Inches(0.4),
             dev, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(10.1), y + Inches(0.5), Inches(2.6), Inches(0.35),
             cmd, 10, MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, Inches(10.1), y + Inches(0.9), Inches(2.6), Inches(0.35),
             "USB / TCP-IP Connected", 10, GREEN, False, PP_ALIGN.CENTER)

# Key Insight box
insight = add_shape(slide, Inches(0.3), Inches(5.5), Inches(12.7), Inches(1.5), BG_CARD2, ACCENT)
add_multi_text(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(1.2),
              [("💡  KEY: Per-User Isolation via X-Client-ID Header", 15, YELLOW, True),
               ("Every browser tab generates a unique clientId (stored in sessionStorage). All API calls include this ID as an HTTP header.", 12, LIGHT_GRAY),
               ("The server maintains a Map<clientId, {serial, adbBinary, model}> so each user independently targets their own device.", 12, LIGHT_GRAY)])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: DLT Bridge — Multi-User Isolation
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "DLT BRIDGE  —  How Each User Gets Isolated Logs", 30, WHITE, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

# Flow diagram: User DLT Viewer → Their Port → Server Proxy → ADB Forward → Device
# Row 1: User A
rows = [
    ("User A", "DLT Viewer\nPort 3490", "TCP Proxy\n0.0.0.0:3490", "ADB Forward\nlocalhost:4490", "Device A\ntcp:3490", ACCENT2, "client_a_1xxx"),
    ("User B", "DLT Viewer\nPort 3491", "TCP Proxy\n0.0.0.0:3491", "ADB Forward\nlocalhost:4491", "Device B\ntcp:3490", GREEN, "client_b_2xxx"),
    ("User C", "DLT Viewer\nPort 3492", "TCP Proxy\n0.0.0.0:3492", "ADB Forward\nlocalhost:4492", "Device A\ntcp:3490", ORANGE, "client_c_3xxx"),
]

# Column headers
headers = ["DLT VIEWER\n(User's PC)", "PUBLIC PORT\n(Unique per user)", "TCP PROXY\n(Node.js net)", "ADB FORWARD\n(Internal port)", "DEVICE\n(Target)"]
for i, h in enumerate(headers):
    x = Inches(0.3 + i * 2.55)
    add_text(slide, x, Inches(1.3), Inches(2.3), Inches(0.7),
             h, 11, MUTED, True, PP_ALIGN.CENTER)

for row_idx, (user, viewer, proxy, fwd, device, color, cid) in enumerate(rows):
    y = Inches(2.2 + row_idx * 1.5)
    
    # User label
    add_shape(slide, Inches(0.3), y, Inches(2.3), Inches(1.1), BG_CARD, color)
    add_text(slide, Inches(0.4), y + Inches(0.05), Inches(2.1), Inches(0.35),
             f"🖥️  {user}", 13, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.4), y + Inches(0.35), Inches(2.1), Inches(0.35),
             viewer, 10, MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.4), y + Inches(0.75), Inches(2.1), Inches(0.25),
             cid, 8, color, False, PP_ALIGN.CENTER)
    
    # Arrow 1
    add_text(slide, Inches(2.55), y + Inches(0.25), Inches(0.4), Inches(0.4),
             "►", 16, color, True, PP_ALIGN.CENTER)
    
    # Proxy
    add_shape(slide, Inches(2.85), y, Inches(2.3), Inches(1.1), BG_CARD, color)
    add_text(slide, Inches(2.95), y + Inches(0.15), Inches(2.1), Inches(0.7),
             proxy, 12, WHITE, True, PP_ALIGN.CENTER)
    
    # Arrow 2
    add_text(slide, Inches(5.1), y + Inches(0.25), Inches(0.4), Inches(0.4),
             "►", 16, color, True, PP_ALIGN.CENTER)
    
    # ADB Forward
    add_shape(slide, Inches(5.4), y, Inches(2.3), Inches(1.1), BG_CARD, color)
    add_text(slide, Inches(5.5), y + Inches(0.15), Inches(2.1), Inches(0.7),
             fwd, 12, WHITE, True, PP_ALIGN.CENTER)
    
    # Arrow 3
    add_text(slide, Inches(7.65), y + Inches(0.25), Inches(0.4), Inches(0.4),
             "►", 16, color, True, PP_ALIGN.CENTER)
    
    # Device
    add_shape(slide, Inches(7.95), y, Inches(2.3), Inches(1.1), BG_CARD, color)
    add_text(slide, Inches(8.05), y + Inches(0.15), Inches(2.1), Inches(0.7),
             device, 12, WHITE, True, PP_ALIGN.CENTER)

# Key insight  
insight_y = Inches(6.0)
add_shape(slide, Inches(0.3), insight_y, Inches(10.0), Inches(1.2), BG_CARD2, ACCENT)
add_multi_text(slide, Inches(0.6), insight_y + Inches(0.1), Inches(9.5), Inches(1.0),
              [("🔑  How Isolation Works", 14, YELLOW, True),
               ("• Each user picks a UNIQUE public port (e.g., 3490, 3491, 3492)", 11, LIGHT_GRAY),
               ("• Server creates an independent TCP proxy per port, tracked by clientId", 11, LIGHT_GRAY),
               ("• When a tab closes (beforeunload), the proxy is auto-destroyed using sendBeacon", 11, LIGHT_GRAY)])

# Cleanup box
add_shape(slide, Inches(10.5), insight_y, Inches(2.5), Inches(1.2), BG_CARD2, RED)
add_multi_text(slide, Inches(10.6), insight_y + Inches(0.1), Inches(2.3), Inches(1.0),
              [("🧹  Auto Cleanup", 13, RED, True),
               ("Tab closes →", 11, LIGHT_GRAY),
               ("Proxy destroyed +", 11, LIGHT_GRAY),
               ("ADB forward removed", 11, LIGHT_GRAY)])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: Global Notifications System
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "GLOBAL NOTIFICATION SYSTEM", 30, WHITE, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "Real-time cross-user notifications — when one user makes a system change, ALL users are alerted instantly.", 15, MUTED)

# How it works
add_shape(slide, Inches(0.4), Inches(1.9), Inches(7.5), Inches(5.0), BG_CARD, RGBColor(0x33, 0x33, 0x44))
add_text(slide, Inches(0.7), Inches(2.0), Inches(7), Inches(0.4),
         "⚙️  How It Works", 18, WHITE, True)

flow_steps = [
    ("1. EVENT TRIGGER", "A user performs an action that affects all users\n(e.g., switches ADB binary, writes IMEI, etc.)", ACCENT2),
    ("2. SERVER STORES", "Server calls addGlobalNotification(type, user, message)\nStored in globalNotifications array (last 10 kept)", GREEN),
    ("3. CLIENT POLLS", "Every 5s, checkDeviceStatus() fetches /api/device-status\nResponse includes notifications[] array", ORANGE),
    ("4. NEW? → POPUP!", "Client checks seenNotifIds (localStorage)\nIf new → shows animated system popup with icon", YELLOW),
]

for i, (step, desc, color) in enumerate(flow_steps):
    y = Inches(2.5 + i * 0.95)
    add_text(slide, Inches(0.7), y, Inches(3), Inches(0.35),
             step, 13, color, True)
    add_text(slide, Inches(3.6), y, Inches(4), Inches(0.8),
             desc, 11, LIGHT_GRAY)

# Notification Types
add_shape(slide, Inches(8.2), Inches(1.9), Inches(4.7), Inches(5.0), BG_CARD, RGBColor(0x33, 0x33, 0x44))
add_text(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(0.4),
         "🔔  Notification Types", 18, WHITE, True)

notif_types = [
    ("⚙️  ADB Binary Change", "When any user switches between\nadb1 / adb2 / custom binary\n→ All clients auto-sync their binary", ACCENT2),
    ("🆔  IMEI Write Event", "When IMEI is written to device\n→ All users see updated IMEI\nimmediately on their dashboard", GREEN),
    ("📢  System Announcements", "Server-push messages for\nmaintenance, updates, or\ncritical device state changes", ORANGE),
]

for i, (title, desc, color) in enumerate(notif_types):
    y = Inches(2.6 + i * 1.45)
    sub_card = add_shape(slide, Inches(8.4), y, Inches(4.3), Inches(1.2), BG_CARD2, color)
    add_text(slide, Inches(8.6), y + Inches(0.05), Inches(4.0), Inches(0.35),
             title, 14, color, True)
    add_text(slide, Inches(8.6), y + Inches(0.4), Inches(4.0), Inches(0.7),
             desc, 11, MUTED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8: Multi-Model Support & Auto-Detection
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "MULTI-MODEL SUPPORT & AUTO-DETECTION", 30, WHITE, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

models = [
    ("🚗  Toyota\n(Single SIM)", "Standard sldd path\n50+ commands\nRegion support: IN, AE, SA, JP", ACCENT2),
    ("🚗  Toyota\n(Dual SIM)", "Dual SIM commands\nSIM slot selection\nExtended telephony", GREEN),
    ("🚗  JLR\n(Single SIM)", "JLR-specific commands\nSingle SIM variant\nCustom validation", ORANGE),
    ("🚗  JLR\n(Dual SIM)", "Dual SIM support\nJLR platform\nFull command set", YELLOW),
    ("🚗  BMW\n(WAVE LOW)", "Path: /usr/bin/factory/sldd\ngetradiostate (numeric)\nsetradiopower 0/1", ACCENT),
    ("🚗  BMW\n(WAVE HIGH)", "Dual SIM BMW\n/usr/bin/factory/sldd\nAuto-detected via SW ver", RED),
]

for i, (name, desc, color) in enumerate(models):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.2)
    y = Inches(1.4 + row * 2.7)
    
    card = add_shape(slide, x, y, Inches(3.9), Inches(2.3), BG_CARD, color)
    add_text(slide, x + Inches(0.15), y + Inches(0.15), Inches(3.6), Inches(0.7),
             name, 16, WHITE, True)
    add_text(slide, x + Inches(0.15), y + Inches(1.0), Inches(3.6), Inches(1.1),
             desc, 12, MUTED)

# Auto-detection insight
insight_y = Inches(6.5)
add_shape(slide, Inches(0.4), insight_y, Inches(12.5), Inches(0.8), BG_CARD2, GREEN)
add_multi_text(slide, Inches(0.7), insight_y + Inches(0.1), Inches(12), Inches(0.6),
              [("🤖  AUTO-DETECTION:  Server reads 'cat etc/version' → If output contains 'WAVE' → auto-switches to BMW model", 14, GREEN, True),
               ("Also uses sticky version cache: if version read fails temporarily, the model doesn't flip back. Smooth experience.", 11, MUTED)])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9: Regression Testing Engine
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "AUTOMATED REGRESSION ENGINE", 30, WHITE, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

# Left side: Features
add_shape(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.8), BG_CARD, RGBColor(0x33, 0x33, 0x44))
add_text(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4),
         "🔄  Regression Features", 20, WHITE, True)

reg_features = [
    "Configure test sequences — drag & drop command ordering",
    "Set iteration count (1 to 1000+ cycles)",
    "Per-step configurable delay (e.g., 5s between commands)",
    "Custom parameters for Dial / SMS steps",
    "Pause / Resume mid-regression without losing progress",
    "Stop at any time — results preserved",
    "Live progress bar with current iteration / total",
    "Real-time PASS / FAIL count on screen",
    "Module-based quick selection (add all SIM, Network, Call commands)",
    "Detailed timestamped log for every step",
]
add_bullet_list(slide, Inches(0.7), Inches(2.1), Inches(5.3), Inches(4.8),
                reg_features, 13, LIGHT_GRAY, GREEN)

# Right side: Export
add_shape(slide, Inches(6.8), Inches(1.4), Inches(6.1), Inches(3.0), BG_CARD, RGBColor(0x33, 0x33, 0x44))
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.4),
         "📊  CSV Export Report", 20, WHITE, True)

add_multi_text(slide, Inches(7.1), Inches(2.1), Inches(5.6), Inches(2.0),
              [("One-click export generates a CSV with:", 13, LIGHT_GRAY),
               ("", 6, MUTED),
               ("• Iteration #  |  Step #  |  Timestamp", 12, ACCENT2),
               ("• Command Name  |  Command ID", 12, ACCENT2),
               ("• Full Output  |  PASS / FAIL Status", 12, ACCENT2),
               ("", 6, MUTED),
               ("File: regression_report_{model}_{timestamp}.csv", 11, MUTED)])

# Smart Validation
add_shape(slide, Inches(6.8), Inches(4.7), Inches(6.1), Inches(2.5), BG_CARD, RGBColor(0x33, 0x33, 0x44))
add_text(slide, Inches(7.1), Inches(4.8), Inches(5.5), Inches(0.4),
         "✅  Smart PASS/FAIL Validation", 20, WHITE, True)

add_multi_text(slide, Inches(7.1), Inches(5.3), Inches(5.6), Inches(1.8),
              [("Commands are auto-validated using regex patterns:", 13, LIGHT_GRAY),
               ("", 6, MUTED),
               ("• SIM State → expects numeric state value", 12, YELLOW),
               ("• Network Type → accepts ANY numeric type (not hardcoded)", 12, YELLOW),
               ("• EID → passes even if empty (eUICC optional)", 12, YELLOW),
               ("• Radio State → supports both Toyota & BMW formats", 12, YELLOW),
               ("• Custom patterns per command ID", 12, YELLOW)])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10: Deployment & Summary
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "DEPLOYMENT & SUMMARY", 30, WHITE, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

# Deployment steps
add_shape(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(3.0), BG_CARD, RGBColor(0x33, 0x33, 0x44))
add_text(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4),
         "🚀  How to Deploy", 20, WHITE, True)

deploy_steps = [
    "1.  Copy TelephonyManager.exe to setup PC",
    "2.  Connect device via USB (or ADB over TCP/IP)",
    "3.  Double-click TelephonyManager.exe to start",
    "4.  Open browser: http://localhost:3000",
    "5.  Share IP with team: http://<PC-IP>:3000",
    "6.  Each user opens their own tab — done! ✅",
]
add_bullet_list(slide, Inches(0.7), Inches(2.1), Inches(5.3), Inches(2.5),
                deploy_steps, 14, LIGHT_GRAY, GREEN)

# Stats / Impact
add_shape(slide, Inches(0.4), Inches(4.7), Inches(6.0), Inches(2.5), BG_CARD, RGBColor(0x33, 0x33, 0x44))
add_text(slide, Inches(0.7), Inches(4.8), Inches(5.5), Inches(0.4),
         "📈  Impact vs Remote PC", 20, WHITE, True)

stats = [
    ("Setup Time", "30+ min → 10 seconds", GREEN),
    ("Concurrent Users", "1 (single session) → Unlimited", GREEN),
    ("Bandwidth Required", "High (screen streaming) → Minimal (JSON API)", GREEN),
    ("Regression Testing", "Fully manual → Fully automated", GREEN),
    ("DLT Setup", "Manual per-PC config → Auto per-user", GREEN),
]

for i, (label, value, color) in enumerate(stats):
    y = Inches(5.3 + i * 0.35)
    add_text(slide, Inches(0.7), y, Inches(2.5), Inches(0.3),
             label, 12, MUTED, True)
    add_text(slide, Inches(3.2), y, Inches(3.0), Inches(0.3),
             value, 12, color, True)

# Thank You
add_shape(slide, Inches(6.8), Inches(1.4), Inches(6.1), Inches(5.8), BG_CARD, ACCENT)
add_text(slide, Inches(7.0), Inches(2.5), Inches(5.7), Inches(1),
         "THANK YOU", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(7.0), Inches(3.5), Inches(5.7), Inches(0.6),
         "Questions & Live Demo", 24, ACCENT2, False, PP_ALIGN.CENTER)

add_multi_text(slide, Inches(7.2), Inches(4.5), Inches(5.3), Inches(2.5),
              [("Developed by:", 14, MUTED),
               ("Nitish Kumar", 18, WHITE, True),
               ("Connected Service Unit  |  LG Electronics", 13, MUTED),
               ("", 10, MUTED),
               ("📧  nitish10.kumar@lge.com", 12, ACCENT2),
               ("🌐  http://<server-ip>:3000", 12, ACCENT2)])

# ─── Save ────────────────────────────────────────────────────────────
output_path = r"D:\OnlineSanity\Telephony_Manager_Demo.pptx"
prs.save(output_path)
print(f"\nPresentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
