"""
Telephony Manager - Demo Presentation Generator (LIGHT THEME)
Generates a professional PPT with light/white theme — same content as dark version
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Light Theme Colors ──────────────────────────────────────────────
BG_WHITE    = RGBColor(0xF8, 0xF9, 0xFC)
BG_CARD     = RGBColor(0xFF, 0xFF, 0xFF)
BG_CARD2    = RGBColor(0xF0, 0xF1, 0xF6)
ACCENT      = RGBColor(0x5B, 0x4C, 0xDB)  # Deep purple
ACCENT2     = RGBColor(0x00, 0x96, 0xC7)  # Professional blue
GREEN       = RGBColor(0x0A, 0x8F, 0x4F)
RED         = RGBColor(0xD1, 0x2F, 0x2F)
ORANGE      = RGBColor(0xE6, 0x7E, 0x00)
TEXT_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_BODY   = RGBColor(0x33, 0x33, 0x44)
TEXT_MUTED  = RGBColor(0x66, 0x66, 0x88)
YELLOW_DARK = RGBColor(0xB8, 0x86, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BORDER      = RGBColor(0xDD, 0xDD, 0xEE)
TITLE_BG    = RGBColor(0x2D, 0x2D, 0x5E)  # Dark header for contrast

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.color.rgb = BORDER
        shape.line.width = Pt(1)
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=TEXT_DARK, bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=TEXT_BODY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Segoe UI'
        p.space_after = Pt(6)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, size=15, color=TEXT_BODY, bullet_color=ACCENT2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run1 = p.add_run()
        run1.text = "●  "
        run1.font.size = Pt(size - 2)
        run1.font.color.rgb = bullet_color
        run1.font.name = 'Segoe UI'
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(size)
        run2.font.color.rgb = color
        run2.font.name = 'Segoe UI'
        p.space_after = Pt(8)
    return txBox


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

# Top accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

# Title block
add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
         "TELEPHONY MANAGER", 52, TITLE_BG, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.8),
         "Online Sanity Testing Platform", 30, ACCENT, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(1.2),
         "A modern web-based tool for remote telephony sanity testing \u2014 replacing heavyweight\n"
         "Remote PC sessions with a lightweight, multi-user, real-time browser experience.",
         16, TEXT_MUTED, False, PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, Inches(0), Inches(6.5), Inches(13.33), Inches(1), BG_CARD2)
add_text(slide, Inches(1), Inches(6.65), Inches(5), Inches(0.6),
         "LG Electronics  |  Connected Service Unit", 13, TEXT_MUTED, False)
add_text(slide, Inches(7), Inches(6.65), Inches(5.5), Inches(0.6),
         "February 2026  |  v1.0", 13, TEXT_MUTED, False, PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: Problem Statement
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "THE PROBLEM  \u2014  Remote PC Limitations", 30, TITLE_BG, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

# Left Card - Remote PC Limitations
card = add_shape(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5), RGBColor(0xFF, 0xF5, 0xF5), RED)
add_text(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
         "Remote PC  (Current Method)", 20, RED, True)

limitations = [
    "Only ONE user can connect at a time (exclusive session)",
    "Requires full desktop streaming \u2014 high bandwidth needed",
    "Slow & laggy UI \u2014 screen refresh delays up to 2-3 seconds",
    "If connection drops, session is lost entirely",
    "No parallel testing \u2014 team waits for one person to finish",
    "Device physical access required to troubleshoot",
    "No automated regression \u2014 everything is manual",
    "DLT logs require separate setup on each PC",
    "No real-time device status dashboard",
    "No test history / export capability",
    "IP conflicts when multiple users try to access"
]
add_bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.2), Inches(4.5),
                limitations, 13, TEXT_BODY, RED)

# Right Card - Our Solution
card = add_shape(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.5), RGBColor(0xF0, 0xFD, 0xF4), GREEN)
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.5),
         "Telephony Manager  (Our Solution)", 20, GREEN, True)

solutions = [
    "Multi-user access \u2014 unlimited tabs, zero conflicts",
    "Lightweight browser UI \u2014 no screen streaming needed",
    "Instant command execution \u2014 sub-second response",
    "Session persists even if browser tab is refreshed",
    "Parallel testing \u2014 team works simultaneously",
    "Remote device access via ADB over network",
    "Built-in automated regression with CSV export",
    "DLT bridge auto-configures per user session",
    "Live device dashboard (IMEI, SIM, Radio, Region)",
    "Full test history with PASS/FAIL export",
    "Per-user isolation \u2014 no IP or port conflicts"
]
add_bullet_list(slide, Inches(7.1), Inches(2.15), Inches(5.5), Inches(4.5),
                solutions, 13, TEXT_BODY, GREEN)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: Key Features
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "KEY FEATURES", 30, TITLE_BG, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

features = [
    ("Live Device Dashboard", "Real-time status bar showing IMEI, SIM State,\nService State, Radio State, Region, SW Version", ACCENT2),
    ("One-Click Commands", "Execute 50+ telephony commands with pass/fail\nvalidation using pattern matching", YELLOW_DARK),
    ("Automated Regression", "Configure test sequences, set iterations, pause/\nresume, and export CSV reports", GREEN),
    ("DLT Bridge", "Auto-configures TCP proxy per user for DLT log\nforwarding \u2014 no manual setup needed", ACCENT),
    ("Multi-User Support", "Each user gets unique client ID, isolated config,\nseparate DLT ports \u2014 zero conflicts", ORANGE),
    ("Global Notifications", "Real-time popups for ADB changes, IMEI writes,\nand system events across all users", RED),
    ("Multi-Model Support", "Toyota, JLR, BMW WAVE LOW/HIGH with auto-\ndetection based on device software version", ACCENT2),
    ("Admin Panel", "Add/remove models, commands, categories\ndynamically without code changes", YELLOW_DARK),
]

# Map accent colors to light card backgrounds
card_bg_map = {
    ACCENT2: RGBColor(0xEF, 0xF8, 0xFF),
    YELLOW_DARK: RGBColor(0xFF, 0xFB, 0xEB),
    GREEN: RGBColor(0xF0, 0xFD, 0xF4),
    ACCENT: RGBColor(0xF3, 0xF0, 0xFF),
    ORANGE: RGBColor(0xFF, 0xF7, 0xED),
    RED: RGBColor(0xFF, 0xF5, 0xF5),
}

icons = ["1", "2", "3", "4", "5", "6", "7", "8"]

for i, (title, desc, color) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.3 + row * 2.9)
    
    cbg = card_bg_map.get(color, BG_CARD)
    card = add_shape(slide, x, y, Inches(2.95), Inches(2.6), cbg, color)
    
    # Number badge
    badge = add_shape(slide, x + Inches(0.12), y + Inches(0.12), Inches(0.4), Inches(0.4), color)
    badge.line.fill.background()
    add_text(slide, x + Inches(0.12), y + Inches(0.12), Inches(0.4), Inches(0.4),
             icons[i], 16, WHITE, True, PP_ALIGN.CENTER)
    
    add_text(slide, x + Inches(0.6), y + Inches(0.15), Inches(2.2), Inches(0.45),
             title, 15, TITLE_BG, True)
    add_text(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.6), Inches(1.8),
             desc, 12, TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: Technology Stack
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "TECHNOLOGY STACK", 30, TITLE_BG, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

tech_stack = [
    ("FRONTEND", [
        ("HTML5 + CSS3 + JavaScript", "Modern responsive UI with dark theme"),
        ("Vanilla JS (No frameworks)", "Zero dependencies, fast load, easy maintenance"),
        ("CSS Variables & Glassmorphism", "Premium dark-mode design system"),
        ("Fetch API + Async/Await", "Non-blocking server communication"),
    ], ACCENT2, RGBColor(0xEF, 0xF8, 0xFF)),
    ("BACKEND", [
        ("Node.js + Express.js", "Lightweight, event-driven server"),
        ("ADB (Android Debug Bridge)", "Device communication via USB/TCP"),
        ("net (TCP Module)", "DLT proxy bridge for log forwarding"),
        ("child_process (execAsync)", "Shell command execution with timeout"),
    ], GREEN, RGBColor(0xF0, 0xFD, 0xF4)),
    ("TOOLING & DEPLOYMENT", [
        ("pkg (by Vercel)", "Compiles to standalone .exe \u2014 no Node.js needed"),
        ("Single Binary Distribution", "TelephonyManager.exe \u2014 just double-click to run"),
        ("sessionStorage + localStorage", "Client-side state persistence per tab"),
        ("JSON-based Config", "Dynamic model/command management via JSON files"),
    ], ORANGE, RGBColor(0xFF, 0xF7, 0xED)),
]

for i, (category, items, color, bg) in enumerate(tech_stack):
    x = Inches(0.4 + i * 4.2)
    card = add_shape(slide, x, Inches(1.4), Inches(4.0), Inches(5.6), bg, color)
    
    # Header bar
    hdr = add_shape(slide, x, Inches(1.4), Inches(4.0), Inches(0.6), color)
    hdr.line.fill.background()
    add_text(slide, x + Inches(0.2), Inches(1.47), Inches(3.5), Inches(0.5),
             category, 18, WHITE, True)
    
    y_offset = 2.2
    for tech, desc in items:
        add_multi_text(slide, x + Inches(0.2), Inches(y_offset), Inches(3.5), Inches(0.9),
                      [(tech, 14, TITLE_BG, True), (desc, 11, TEXT_MUTED)])
        y_offset += 0.85


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: Architecture
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "ARCHITECTURE  \u2014  How the Server Connects", 30, TITLE_BG, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

# --- Users Column ---
add_text(slide, Inches(0.3), Inches(1.3), Inches(3), Inches(0.4),
         "BROWSER CLIENTS", 16, ACCENT2, True, PP_ALIGN.CENTER)

users = ["User A  (Tab 1)", "User B  (Tab 2)", "User C  (Tab 3)"]
for i, user in enumerate(users):
    y = Inches(1.9 + i * 1.2)
    card = add_shape(slide, Inches(0.3), y, Inches(3.0), Inches(0.9), RGBColor(0xEF, 0xF8, 0xFF), ACCENT2)
    add_text(slide, Inches(0.5), y + Inches(0.05), Inches(2.5), Inches(0.4),
             user, 14, TITLE_BG, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), y + Inches(0.45), Inches(2.5), Inches(0.35),
             f"clientId: client_{'abc'[i]}_{i+1}xxx", 10, TEXT_MUTED, False, PP_ALIGN.CENTER)

# Arrows
for i in range(3):
    y = Inches(2.2 + i * 1.2)
    add_text(slide, Inches(3.3), y, Inches(0.8), Inches(0.4),
             "\u2500\u2500\u25ba", 18, ACCENT2, True, PP_ALIGN.CENTER)

# --- Server Column ---
add_text(slide, Inches(4.1), Inches(1.3), Inches(5), Inches(0.4),
         "EXPRESS.JS SERVER  (Port 3000)", 16, GREEN, True, PP_ALIGN.CENTER)

server_card = add_shape(slide, Inches(4.1), Inches(1.9), Inches(5.0), Inches(4.8), RGBColor(0xF0, 0xFD, 0xF4), GREEN)

server_items = [
    ("Request Router", "Maps X-Client-ID header to per-user config\n(userConfigs Map: clientId -> {serial, adb, model})", TITLE_BG),
    ("Command Engine", "Executes ADB commands via child_process\nwith timeout, caching & pattern validation", TITLE_BG),
    ("Device Cache", "10s TTL cache with Promise.allSettled\nStale-cache fallback prevents UI flickering", TITLE_BG),
    ("Notification Bus", "Global notification array polled by all clients\nADB changes, IMEI writes broadcast to everyone", TITLE_BG),
]

for i, (title, desc, color) in enumerate(server_items):
    y = Inches(2.0 + i * 1.1)
    add_text(slide, Inches(4.3), y, Inches(4.5), Inches(0.35),
             f"{i+1}. {title}", 13, color, True)
    add_text(slide, Inches(4.3), y + Inches(0.3), Inches(4.5), Inches(0.6),
             desc, 10, TEXT_MUTED)

# Arrows
for i in range(2):
    y = Inches(2.2 + i * 2.0)
    add_text(slide, Inches(9.1), y, Inches(0.8), Inches(0.4),
             "\u2500\u2500\u25ba", 18, GREEN, True, PP_ALIGN.CENTER)

# --- Devices Column ---
add_text(slide, Inches(9.9), Inches(1.3), Inches(3), Inches(0.4),
         "TARGET DEVICES", 16, ORANGE, True, PP_ALIGN.CENTER)

devices = [
    ("Device A (Toyota)", "adb1 -s SERIAL_A"),
    ("Device B (BMW)", "adb1 -s SERIAL_B"),
]
for i, (dev, cmd) in enumerate(devices):
    y = Inches(1.9 + i * 2.0)
    card = add_shape(slide, Inches(9.9), y, Inches(3.0), Inches(1.4), RGBColor(0xFF, 0xF7, 0xED), ORANGE)
    add_text(slide, Inches(10.1), y + Inches(0.1), Inches(2.6), Inches(0.4),
             dev, 14, TITLE_BG, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(10.1), y + Inches(0.5), Inches(2.6), Inches(0.35),
             cmd, 10, TEXT_MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, Inches(10.1), y + Inches(0.9), Inches(2.6), Inches(0.35),
             "USB / TCP-IP Connected", 10, GREEN, False, PP_ALIGN.CENTER)

# Key Insight
insight = add_shape(slide, Inches(0.3), Inches(5.5), Inches(12.7), Inches(1.5), RGBColor(0xF3, 0xF0, 0xFF), ACCENT)
add_multi_text(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(1.2),
              [("KEY: Per-User Isolation via X-Client-ID Header", 15, ACCENT, True),
               ("Every browser tab generates a unique clientId (stored in sessionStorage). All API calls include this ID as an HTTP header.", 12, TEXT_BODY),
               ("The server maintains a Map<clientId, {serial, adbBinary, model}> so each user independently targets their own device.", 12, TEXT_BODY)])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: DLT Bridge
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "DLT BRIDGE  \u2014  How Each User Gets Isolated Logs", 30, TITLE_BG, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

# Column headers
headers = ["DLT VIEWER\n(User's PC)", "PUBLIC PORT\n(Unique per user)", "TCP PROXY\n(Node.js net)", "ADB FORWARD\n(Internal port)", "DEVICE\n(Target)"]
for i, h in enumerate(headers):
    x = Inches(0.3 + i * 2.55)
    add_text(slide, x, Inches(1.3), Inches(2.3), Inches(0.7),
             h, 11, TEXT_MUTED, True, PP_ALIGN.CENTER)

rows = [
    ("User A", "DLT Viewer\nPort 3490", "TCP Proxy\n0.0.0.0:3490", "ADB Forward\nlocalhost:4490", "Device A\ntcp:3490", ACCENT2, "client_a_1xxx", RGBColor(0xEF, 0xF8, 0xFF)),
    ("User B", "DLT Viewer\nPort 3491", "TCP Proxy\n0.0.0.0:3491", "ADB Forward\nlocalhost:4491", "Device B\ntcp:3490", GREEN, "client_b_2xxx", RGBColor(0xF0, 0xFD, 0xF4)),
    ("User C", "DLT Viewer\nPort 3492", "TCP Proxy\n0.0.0.0:3492", "ADB Forward\nlocalhost:4492", "Device A\ntcp:3490", ORANGE, "client_c_3xxx", RGBColor(0xFF, 0xF7, 0xED)),
]

for row_idx, (user, viewer, proxy, fwd, device, color, cid, bg) in enumerate(rows):
    y = Inches(2.2 + row_idx * 1.5)
    
    # User
    add_shape(slide, Inches(0.3), y, Inches(2.3), Inches(1.1), bg, color)
    add_text(slide, Inches(0.4), y + Inches(0.05), Inches(2.1), Inches(0.35),
             user, 13, TITLE_BG, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.4), y + Inches(0.35), Inches(2.1), Inches(0.35),
             viewer, 10, TEXT_MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.4), y + Inches(0.75), Inches(2.1), Inches(0.25),
             cid, 8, color, False, PP_ALIGN.CENTER)
    
    add_text(slide, Inches(2.55), y + Inches(0.25), Inches(0.4), Inches(0.4),
             "\u25ba", 16, color, True, PP_ALIGN.CENTER)
    
    # Proxy
    add_shape(slide, Inches(2.85), y, Inches(2.3), Inches(1.1), bg, color)
    add_text(slide, Inches(2.95), y + Inches(0.15), Inches(2.1), Inches(0.7),
             proxy, 12, TITLE_BG, True, PP_ALIGN.CENTER)
    
    add_text(slide, Inches(5.1), y + Inches(0.25), Inches(0.4), Inches(0.4),
             "\u25ba", 16, color, True, PP_ALIGN.CENTER)
    
    # ADB Forward
    add_shape(slide, Inches(5.4), y, Inches(2.3), Inches(1.1), bg, color)
    add_text(slide, Inches(5.5), y + Inches(0.15), Inches(2.1), Inches(0.7),
             fwd, 12, TITLE_BG, True, PP_ALIGN.CENTER)
    
    add_text(slide, Inches(7.65), y + Inches(0.25), Inches(0.4), Inches(0.4),
             "\u25ba", 16, color, True, PP_ALIGN.CENTER)
    
    # Device
    add_shape(slide, Inches(7.95), y, Inches(2.3), Inches(1.1), bg, color)
    add_text(slide, Inches(8.05), y + Inches(0.15), Inches(2.1), Inches(0.7),
             device, 12, TITLE_BG, True, PP_ALIGN.CENTER)

# Key insight
insight_y = Inches(6.0)
add_shape(slide, Inches(0.3), insight_y, Inches(10.0), Inches(1.2), RGBColor(0xF3, 0xF0, 0xFF), ACCENT)
add_multi_text(slide, Inches(0.6), insight_y + Inches(0.1), Inches(9.5), Inches(1.0),
              [("How Isolation Works", 14, ACCENT, True),
               ("Each user picks a UNIQUE public port (e.g., 3490, 3491, 3492)", 11, TEXT_BODY),
               ("Server creates an independent TCP proxy per port, tracked by clientId", 11, TEXT_BODY),
               ("When a tab closes (beforeunload), the proxy is auto-destroyed using sendBeacon", 11, TEXT_BODY)])

# Cleanup box
add_shape(slide, Inches(10.5), insight_y, Inches(2.5), Inches(1.2), RGBColor(0xFF, 0xF5, 0xF5), RED)
add_multi_text(slide, Inches(10.6), insight_y + Inches(0.1), Inches(2.3), Inches(1.0),
              [("Auto Cleanup", 13, RED, True),
               ("Tab closes ->", 11, TEXT_BODY),
               ("Proxy destroyed +", 11, TEXT_BODY),
               ("ADB forward removed", 11, TEXT_BODY)])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: Global Notifications
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "GLOBAL NOTIFICATION SYSTEM", 30, TITLE_BG, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "Real-time cross-user notifications \u2014 when one user makes a system change, ALL users are alerted instantly.", 15, TEXT_MUTED)

# How it works
add_shape(slide, Inches(0.4), Inches(1.9), Inches(7.5), Inches(5.0), BG_CARD, BORDER)
add_text(slide, Inches(0.7), Inches(2.0), Inches(7), Inches(0.4),
         "How It Works", 18, TITLE_BG, True)

flow_steps = [
    ("1. EVENT TRIGGER", "A user performs an action that affects all users\n(e.g., switches ADB binary, writes IMEI, etc.)", ACCENT2),
    ("2. SERVER STORES", "Server calls addGlobalNotification(type, user, message)\nStored in globalNotifications array (last 10 kept)", GREEN),
    ("3. CLIENT POLLS", "Every 5s, checkDeviceStatus() fetches /api/device-status\nResponse includes notifications[] array", ORANGE),
    ("4. NEW? -> POPUP!", "Client checks seenNotifIds (localStorage)\nIf new -> shows animated system popup with icon", ACCENT),
]

for i, (step, desc, color) in enumerate(flow_steps):
    y = Inches(2.5 + i * 0.95)
    add_text(slide, Inches(0.7), y, Inches(3), Inches(0.35),
             step, 13, color, True)
    add_text(slide, Inches(3.6), y, Inches(4), Inches(0.8),
             desc, 11, TEXT_BODY)

# Notification Types
add_shape(slide, Inches(8.2), Inches(1.9), Inches(4.7), Inches(5.0), BG_CARD, BORDER)
add_text(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(0.4),
         "Notification Types", 18, TITLE_BG, True)

notif_types = [
    ("ADB Binary Change", "When any user switches between\nadb1 / adb2 / custom binary\n-> All clients auto-sync their binary", ACCENT2, RGBColor(0xEF, 0xF8, 0xFF)),
    ("IMEI Write Event", "When IMEI is written to device\n-> All users see updated IMEI\nimmediately on their dashboard", GREEN, RGBColor(0xF0, 0xFD, 0xF4)),
    ("System Announcements", "Server-push messages for\nmaintenance, updates, or\ncritical device state changes", ORANGE, RGBColor(0xFF, 0xF7, 0xED)),
]

for i, (title, desc, color, bg) in enumerate(notif_types):
    y = Inches(2.6 + i * 1.45)
    sub_card = add_shape(slide, Inches(8.4), y, Inches(4.3), Inches(1.2), bg, color)
    add_text(slide, Inches(8.6), y + Inches(0.05), Inches(4.0), Inches(0.35),
             title, 14, color, True)
    add_text(slide, Inches(8.6), y + Inches(0.4), Inches(4.0), Inches(0.7),
             desc, 11, TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8: Multi-Model Support
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "MULTI-MODEL SUPPORT & AUTO-DETECTION", 30, TITLE_BG, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

models = [
    ("Toyota\n(Single SIM)", "Standard sldd path\n50+ commands\nRegion support: IN, AE, SA, JP", ACCENT2, RGBColor(0xEF, 0xF8, 0xFF)),
    ("Toyota\n(Dual SIM)", "Dual SIM commands\nSIM slot selection\nExtended telephony", GREEN, RGBColor(0xF0, 0xFD, 0xF4)),
    ("JLR\n(Single SIM)", "JLR-specific commands\nSingle SIM variant\nCustom validation", ORANGE, RGBColor(0xFF, 0xF7, 0xED)),
    ("JLR\n(Dual SIM)", "Dual SIM support\nJLR platform\nFull command set", ACCENT, RGBColor(0xF3, 0xF0, 0xFF)),
    ("BMW\n(WAVE LOW)", "Path: /usr/bin/factory/sldd\ngetradiostate (numeric)\nsetradiopower 0/1", RED, RGBColor(0xFF, 0xF5, 0xF5)),
    ("BMW\n(WAVE HIGH)", "Dual SIM BMW\n/usr/bin/factory/sldd\nAuto-detected via SW ver", YELLOW_DARK, RGBColor(0xFF, 0xFB, 0xEB)),
]

for i, (name, desc, color, bg) in enumerate(models):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.2)
    y = Inches(1.4 + row * 2.7)
    
    card = add_shape(slide, x, y, Inches(3.9), Inches(2.3), bg, color)
    add_text(slide, x + Inches(0.15), y + Inches(0.15), Inches(3.6), Inches(0.7),
             name, 16, TITLE_BG, True)
    add_text(slide, x + Inches(0.15), y + Inches(1.0), Inches(3.6), Inches(1.1),
             desc, 12, TEXT_MUTED)

# Auto-detection
insight_y = Inches(6.5)
add_shape(slide, Inches(0.4), insight_y, Inches(12.5), Inches(0.8), RGBColor(0xF0, 0xFD, 0xF4), GREEN)
add_multi_text(slide, Inches(0.7), insight_y + Inches(0.1), Inches(12), Inches(0.6),
              [("AUTO-DETECTION:  Server reads 'cat etc/version'  ->  If output contains 'WAVE'  ->  auto-switches to BMW model", 14, GREEN, True),
               ("Also uses sticky version cache: if version read fails temporarily, the model doesn't flip back. Smooth experience.", 11, TEXT_MUTED)])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9: Regression Engine
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "AUTOMATED REGRESSION ENGINE", 30, TITLE_BG, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

# Left: Features
add_shape(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.8), BG_CARD, BORDER)
add_text(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4),
         "Regression Features", 20, TITLE_BG, True)

reg_features = [
    "Configure test sequences \u2014 drag & drop command ordering",
    "Set iteration count (1 to 1000+ cycles)",
    "Per-step configurable delay (e.g., 5s between commands)",
    "Custom parameters for Dial / SMS steps",
    "Pause / Resume mid-regression without losing progress",
    "Stop at any time \u2014 results preserved",
    "Live progress bar with current iteration / total",
    "Real-time PASS / FAIL count on screen",
    "Module-based quick selection (add all SIM, Network, Call commands)",
    "Detailed timestamped log for every step",
]
add_bullet_list(slide, Inches(0.7), Inches(2.1), Inches(5.3), Inches(4.8),
                reg_features, 13, TEXT_BODY, GREEN)

# Right top: Export
add_shape(slide, Inches(6.8), Inches(1.4), Inches(6.1), Inches(3.0), BG_CARD, BORDER)
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.4),
         "CSV Export Report", 20, TITLE_BG, True)

add_multi_text(slide, Inches(7.1), Inches(2.1), Inches(5.6), Inches(2.0),
              [("One-click export generates a CSV with:", 13, TEXT_BODY),
               ("", 6, TEXT_MUTED),
               ("Iteration #  |  Step #  |  Timestamp", 12, ACCENT2),
               ("Command Name  |  Command ID", 12, ACCENT2),
               ("Full Output  |  PASS / FAIL Status", 12, ACCENT2),
               ("", 6, TEXT_MUTED),
               ("File: regression_report_{model}_{timestamp}.csv", 11, TEXT_MUTED)])

# Right bottom: Validation
add_shape(slide, Inches(6.8), Inches(4.7), Inches(6.1), Inches(2.5), BG_CARD, BORDER)
add_text(slide, Inches(7.1), Inches(4.8), Inches(5.5), Inches(0.4),
         "Smart PASS/FAIL Validation", 20, TITLE_BG, True)

add_multi_text(slide, Inches(7.1), Inches(5.3), Inches(5.6), Inches(1.8),
              [("Commands are auto-validated using regex patterns:", 13, TEXT_BODY),
               ("", 6, TEXT_MUTED),
               ("SIM State  ->  expects numeric state value", 12, YELLOW_DARK),
               ("Network Type  ->  accepts ANY numeric type (not hardcoded)", 12, YELLOW_DARK),
               ("EID  ->  passes even if empty (eUICC optional)", 12, YELLOW_DARK),
               ("Radio State  ->  supports both Toyota & BMW formats", 12, YELLOW_DARK),
               ("Custom patterns per command ID", 12, YELLOW_DARK)])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10: Deployment & Summary
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "DEPLOYMENT & SUMMARY", 30, TITLE_BG, True)
add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), ACCENT)

# Deployment steps
add_shape(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(3.0), BG_CARD, BORDER)
add_text(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4),
         "How to Deploy", 20, TITLE_BG, True)

deploy_steps = [
    "1.  Copy TelephonyManager.exe to setup PC",
    "2.  Connect device via USB (or ADB over TCP/IP)",
    "3.  Double-click TelephonyManager.exe to start",
    "4.  Open browser: http://localhost:3000",
    "5.  Share IP with team: http://<PC-IP>:3000",
    "6.  Each user opens their own tab \u2014 done!",
]
add_bullet_list(slide, Inches(0.7), Inches(2.1), Inches(5.3), Inches(2.5),
                deploy_steps, 14, TEXT_BODY, GREEN)

# Stats
add_shape(slide, Inches(0.4), Inches(4.7), Inches(6.0), Inches(2.5), BG_CARD, BORDER)
add_text(slide, Inches(0.7), Inches(4.8), Inches(5.5), Inches(0.4),
         "Impact vs Remote PC", 20, TITLE_BG, True)

stats = [
    ("Setup Time", "30+ min  ->  10 seconds", GREEN),
    ("Concurrent Users", "1 (single session)  ->  Unlimited", GREEN),
    ("Bandwidth Required", "High (screen streaming)  ->  Minimal (JSON API)", GREEN),
    ("Regression Testing", "Fully manual  ->  Fully automated", GREEN),
    ("DLT Setup", "Manual per-PC config  ->  Auto per-user", GREEN),
]

for i, (label, value, color) in enumerate(stats):
    y = Inches(5.3 + i * 0.35)
    add_text(slide, Inches(0.7), y, Inches(2.5), Inches(0.3),
             label, 12, TEXT_MUTED, True)
    add_text(slide, Inches(3.2), y, Inches(3.0), Inches(0.3),
             value, 12, color, True)

# Thank You
add_shape(slide, Inches(6.8), Inches(1.4), Inches(6.1), Inches(5.8), ACCENT, ACCENT)
add_text(slide, Inches(7.0), Inches(2.5), Inches(5.7), Inches(1),
         "THANK YOU", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(7.0), Inches(3.5), Inches(5.7), Inches(0.6),
         "Questions & Live Demo", 24, RGBColor(0xCC, 0xCC, 0xFF), False, PP_ALIGN.CENTER)

add_multi_text(slide, Inches(7.2), Inches(4.5), Inches(5.3), Inches(2.5),
              [("Developed by:", 14, RGBColor(0xCC, 0xCC, 0xEE)),
               ("Nitish Kumar", 18, WHITE, True),
               ("Connected Service Unit  |  LG Electronics", 13, RGBColor(0xCC, 0xCC, 0xEE)),
               ("", 10, WHITE),
               ("nitish10.kumar@lge.com", 12, RGBColor(0xCC, 0xDD, 0xFF)),
               ("http://<server-ip>:3000", 12, RGBColor(0xCC, 0xDD, 0xFF))])

# ─── Save ────────────────────────────────────────────────────────────
output_path = r"D:\OnlineSanity\Telephony_Manager_Demo_Light.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
